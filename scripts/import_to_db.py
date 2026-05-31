"""
Импорт из Excel в PostgreSQL: условие (текст) + ссылка на файл + извлечённый текст файла.

Зависимости (в активном venv проекта):
    pip install pandas openpyxl psycopg2-binary requests \\
                pdfplumber python-docx python-pptx beautifulsoup4 lxml

Создание таблицы (один раз):
    psql -U postgres -d itflux -f scripts/sql/fipi_part2_import.sql

Запуск:
    python scripts/import_to_db.py
    # или с переопределением пути к Excel:
    EXCEL_FILE="/path/to/file.xlsx" python scripts/import_to_db.py
"""

import io
import os
import sys
import zipfile
from urllib.parse import urlparse

import pandas as pd
import psycopg2
import requests

# ─────────────────────────────────────────────
#  КОНФИГУРАЦИЯ
# ─────────────────────────────────────────────

# Путь к Excel (шапка: text_plain | files)
_EXCEL_DEFAULT = (
    "/Users/darsorokina/Desktop/уроки/Информатика/задачи части 2.xlsx"
)
EXCEL_FILE = os.environ.get("EXCEL_FILE", _EXCEL_DEFAULT)

# Имена колонок в Excel (как в первой строке файла)
COL_CONDITION = "text_plain"
COL_FILE_URL = "files"

# PostgreSQL — те же переменные, что и Django (Generator.settings)
DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "itflux"),
    "user": os.environ.get("PGUSER", "postgres"),
    "password": os.environ.get("PGPASSWORD", "postgres"),
}

TABLE_NAME = os.environ.get("IMPORT_TABLE", "fipi_part2_import")
DB_COL_CONDITION = "task_condition"
DB_COL_FILE_TEXT = "file_content"
DB_COL_FILE_URL = "file_url"

SKIP_ON_ERROR = os.environ.get("SKIP_ON_ERROR", "").lower() in ("1", "true", "yes")

# ─────────────────────────────────────────────


def extract_text_pdf(data: bytes) -> str:
    import pdfplumber

    text_parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    return "\n".join(text_parts)


def extract_text_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def extract_text_xlsx(data: bytes) -> str:
    df = pd.read_excel(io.BytesIO(data), sheet_name=None)
    parts = []
    for sheet_name, sheet_df in df.items():
        parts.append(f"[Sheet: {sheet_name}]")
        parts.append(sheet_df.to_string(index=False))
    return "\n".join(parts)


def extract_text_html(data: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data, "lxml")
    return soup.get_text(separator="\n", strip=True)


def extract_text_zip(data: bytes) -> str:
    parts = []
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        for name in zf.namelist():
            if name.endswith("/") or os.path.basename(name).startswith((".", "__")):
                continue
            try:
                file_data = zf.read(name)
                text = extract_text(file_data, name)
                parts.append(f"[{name}]\n{text}")
                print(f"             zip → {name}: {len(text)} символов")
            except Exception as e:
                parts.append(f"[{name}]\n[Ошибка извлечения: {e}]")
    return "\n\n".join(parts)


def extract_text(data: bytes, filename: str) -> str:
    path = urlparse(filename).path.lower() if filename.startswith("http") else filename.lower()

    if path.endswith(".zip"):
        return extract_text_zip(data)
    if path.endswith(".pdf"):
        return extract_text_pdf(data)
    if path.endswith(".docx"):
        return extract_text_docx(data)
    if path.endswith(".doc"):
        try:
            return extract_text_docx(data)
        except Exception:
            return data.decode("utf-8", errors="replace")
    if path.endswith(".pptx"):
        return extract_text_pptx(data)
    if path.endswith((".xlsx", ".xls")):
        return extract_text_xlsx(data)
    if path.endswith((".html", ".htm")):
        return extract_text_html(data)
    if path.endswith((".txt", ".md", ".csv")):
        return data.decode("utf-8", errors="replace")
    try:
        return data.decode("utf-8")
    except Exception:
        return f"[Не удалось извлечь текст: неизвестный формат ({path})]"


def download_file(url: str, timeout: int = 60) -> bytes:
    resp = requests.get(url, timeout=timeout)
    resp.raise_for_status()
    return resp.content


def main():
    if not os.path.isfile(EXCEL_FILE):
        sys.exit(
            f"Файл не найден: {EXCEL_FILE}\n"
            "Укажите путь: EXCEL_FILE=/полный/путь.xlsx python scripts/import_to_db.py"
        )

    print(f"Читаю файл: {EXCEL_FILE}")
    df = pd.read_excel(EXCEL_FILE)

    if COL_CONDITION not in df.columns:
        sys.exit(f"Нет колонки '{COL_CONDITION}'. Есть: {list(df.columns)}")
    if COL_FILE_URL not in df.columns:
        sys.exit(f"Нет колонки '{COL_FILE_URL}'. Есть: {list(df.columns)}")

    total = len(df)
    print(f"Строк: {total}\n")

    print("Подключение к PostgreSQL...")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    inserted = 0
    skipped = 0
    errors = 0

    for idx, row in df.iterrows():
        condition = str(row[COL_CONDITION]).strip() if pd.notna(row[COL_CONDITION]) else ""
        url = str(row[COL_FILE_URL]).strip() if pd.notna(row[COL_FILE_URL]) else ""

        short_url = (url[:80] + "…") if len(url) > 80 else url
        print(f"[{idx + 1}/{total}] URL: {short_url or '(пусто)'}")

        file_text = None

        if url:
            try:
                data = download_file(url)
                file_text = extract_text(data, url)
                print(f"         ✓ извлечено {len(file_text)} символов")
            except Exception as e:
                errors += 1
                print(f"         ✗ ошибка: {e}")
                if SKIP_ON_ERROR:
                    skipped += 1
                    print("         → строка пропущена")
                    continue
        else:
            print("         ! URL пустой")

        try:
            cur.execute(
                f"""
                INSERT INTO {TABLE_NAME} ({DB_COL_CONDITION}, {DB_COL_FILE_TEXT}, {DB_COL_FILE_URL})
                VALUES (%s, %s, %s)
                """,
                (condition, file_text, url),
            )
            conn.commit()
            inserted += 1
        except Exception as e:
            conn.rollback()
            errors += 1
            print(f"         ✗ ошибка вставки в БД: {e}")
            if SKIP_ON_ERROR:
                skipped += 1

    cur.close()
    conn.close()

    print("\n─────────────────────────────")
    print(f"Готово! Вставлено: {inserted}, пропущено: {skipped}, ошибок: {errors}")


if __name__ == "__main__":
    main()
